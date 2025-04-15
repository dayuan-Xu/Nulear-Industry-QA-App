# 该模块基于LangChain的Document loaders：https://python.langchain.com/docs/how_to/#document-loaders
# 各个函数加载指定路径下的文件，返回一个List[Document]
# 对于文件的加载要求：仅仅要求提取出文件中的文本内容即可，尽量将文件中的图表、表格等复杂结构也转化为文本。
# 要求各个函数将文件内容加载为List[Document]后输出前几个doc，查看是否正确加载文件内容为文本。

import os
from dotenv import load_dotenv
from typing_extensions import List
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_unstructured import UnstructuredLoader
from langchain_community.document_loaders import UnstructuredMarkdownLoader
from langchain_community.document_loaders import Docx2txtLoader
from pptx import Presentation
from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader

# 加载环境变量
load_dotenv(override=True)
UNSTRUCTURED_API_KEY=os.getenv('UNSTRUCTURED_API_KEY')
AZURE_ENDPOINT=os.getenv('AZURE_ENDPOINT')
AZURE_API_KEY=os.getenv('AZURE_API_KEY')
print("付费服务相关变量设置成功")
def load_txt(file_path:str)->List[Document]:
    # 该函数加载一个TXT文件，返回List[Document]

    # 1、建立文本切分器
    text_splitter = RecursiveCharacterTextSplitter(
        separators=[
            "\n\n",  # 空行，常作为不同章节的分割标志
            "\n",  # 换行符，常作为段落之间或文本与公式之间的分割标志
            " ",  # 空格，常作为英文单词之间分割标志
            ".",
            ",",
            "\u200b",  # Zero-width space (零宽空格): （不可见）
            "\uff0c",  # Fullwidth comma (全角逗号): ，
            "\u3001",  # Ideographic comma (顿号): 、
            "\uff0e",  # Fullwidth full stop (全角句号): ．
            "\u3002",  # Ideographic full stop (句号): 。
            "",
        ],
        chunk_size=300, chunk_overlap=50, add_start_index=True
    )
    # 2、获取所有文档对象
    all_docs = []
    # 读文件，显式指定编码为 utf-8
    try:
        with open(file_path, encoding='utf-8') as f:
            content = ""
            while True:
                chunk = f.read(10000)  # 每次读取10000字节
                if not chunk:
                    break
                content += chunk
                # 每读取一定量的内容进行切分
                if len(content) > 100000:  # 根据需要调整阈值
                    texts = text_splitter.create_documents([content])
                    all_docs.extend(texts)
                    content = ""
    except UnicodeDecodeError:
        print(f"文件 {file_path} 使用的编码不是 utf-8，尝试使用 GBK 编码...")
        with open(file_path, encoding='gbk', errors='replace') as f:
            content = ""
            while True:
                chunk = f.read(10000)  # 每次读取10000字节
                if not chunk:
                    break
                content += chunk
                # 每读取一定量的内容进行切分
                if len(content) > 100000:  # 根据需要调整阈值
                    texts = text_splitter.create_documents([content])
                    all_docs.extend(texts)
                    content = ""

    # 处理剩余的内容
    if content:
        texts = text_splitter.create_documents([content])
        all_docs.extend(texts)

    # 3、为每一个doc的元数据添加一个属性metadat["source"]=file_path
    for i in range(len(all_docs)):
        all_docs[i].metadata["source"]=file_path
    print("测试用的TXT文件的被切分成为:", len(all_docs), "份\n")
    # 输出前三个元素查看是否成功加载
    for i in range(3):
        print(all_docs[i])
    # 4、返回该List[Document]
    return all_docs
def load_pdf_simply(file_path:str)->List[Document]:
    # 该方法基于pypdf库加载pdf文件，实现对pdf文件中文本的简单快速提取。

    # 1、初始化加载器
    loader = PyPDFLoader(file_path=file_path)
    # 2、获取Doc列表，类型为List[Document]，此时每个doc对应pdf文件中的一页，显然需要对其内容进行切分
    docs = loader.load()
    # 3、输出该pdf文件中的总页数
    print("该pdf文件总页数：",len(docs))
    # 4、输出查看该pdf文件第一页的元数据和具体内容
    for key, value in docs[0].metadata.items():#遍历docs[0].metadata中的每一个键值对
        print(f"{key}: {value}")
    print(docs[0].page_content)
    # 5、对pdf文件中的每一页进行切分，并返回一个List[Document]
    text_splitter = RecursiveCharacterTextSplitter(
        separators=[
            "\n\n",  # 空行，常作为不同章节的分割标志
            "\n",  # 换行符，常作为段落之间或文本与公式之间的分割标志
            " ",  # 空格，常作为英文单词之间分割标志
            ".",
            ",",
            "\u200b",  # Zero-width space (零宽空格): （不可见）
            "\uff0c",  # Fullwidth comma (全角逗号): ，
            "\u3001",  # Ideographic comma (顿号): 、
            "\uff0e",  # Fullwidth full stop (全角句号): ．
            "\u3002",  # Ideographic full stop (句号): 。
            "",
        ],
        chunk_size=500, chunk_overlap=200, add_start_index=True
    )
    all_splits=text_splitter.split_documents(docs)
    return all_splits

def load_pdf_with_Unstructured(file_path:str):
    # 该方法测试pdf文件的布局解析，使用Unstructured提供的API接口。
    # 优点：能够正确识别pdf文件中不同的结构类别，从而可以编码实现专门从指定结构中提取文本: 比如可以从表格中正确提取数据，从富文本图像中正确提取文本。
    # 缺点：对于公式、两行内容的表头识别可能不准确（可能与pdf文件本身清晰度有关）；没有把跨左右两栏的一个段落识别为一个整体

    # 1、初始化加载器
    if UNSTRUCTURED_API_KEY is None:
        print("未设置UNSTRUCTURED_API_KEY，请先加载该变量！")
        return  []
    loader = UnstructuredLoader(
        file_path=file_path,
        strategy="hi_res",
        partition_via_api=True,
        api_key=UNSTRUCTURED_API_KEY,
    )

    # 2、获取Doc列表
    docs = []
    # 此时其中每个Document对应PDF页面中的一个结构类别，
    # 结构类别包括：标题Title（节标题或者图表标题）、页眉Header、页码PageNumber、表格Table、图像Image、
    # 叙述性文本NarrativeText（正文段落或者图表名）、公式Formula以及无法识别类别的UncategorizedText
    for doc in loader.lazy_load():
        docs.append(doc)

    # 3、输出该pdf文件中的结构总数
    print("该pdf文件中的结构总数: ",len(docs),end="\n")

    # 4、输出查看指定页的所有结构
    page1_docs = [doc for doc in docs if doc.metadata.get("page_number") == 3]
    for doc in page1_docs:
        for key,value in doc.metadata.items():
            if key in ["parent_id","category","element_id","text_as_html"]:
                print(f"{key}: {value}")
        print("page_content:",doc.page_content)
        print()

def load_md(file_path:str)->List[Document]:
    # 负责人：么一明
    # 该方法加载md文件，并返回一个List[Document]
    try:
        # 使用 LangChain 提供的 UnstructuredMarkdownLoader 加载 Markdown 文件
        loader = UnstructuredMarkdownLoader(file_path)
        # 加载文档并返回
        documents = loader.load()
        # 初始化 RecursiveCharacterTextSplitter
        text_splitter = RecursiveCharacterTextSplitter(
            separators=[
                "\n\n",  # 空行，常作为不同章节的分割标志
                "\n",  # 换行符，常作为段落之间或文本与公式之间的分割标志
                " ",  # 空格，常作为英文单词之间分割标志
                ".",
                ",",
                "\u200b",  # Zero-width space (零宽空格): （不可见）
                "\uff0c",  # Fullwidth comma (全角逗号): ，
                "\u3001",  # Ideographic comma (顿号): 、
                "\uff0e",  # Fullwidth full stop (全角句号): ．
                "\u3002",  # Ideographic full stop (句号): 。
            ],
            chunk_size=500,
            chunk_overlap=200,
            add_start_index=True
        )
        # 分割文档
        all_splits = text_splitter.split_documents(documents)
        print(all_splits[0].metadata,end="\n")
        print(all_splits[0].page_content)
        return all_splits
    except Exception as e:
        print(f"加载 Markdown 文件时出错: {e}")
        return []

def load_docx_simply(file_path:str)->List[Document]:
    # 负责人：李宏伟
    # 该方法加载docx文件，并返回一个List[Document]
    # 1. 检查文件是否存在
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件 {file_path} 不存在！")

    # 2. 初始化加载器
    loader = Docx2txtLoader(file_path)

    # 3. 加载文档
    docs = loader.load()

    # 4. 初始化文本分割器
    text_splitter = RecursiveCharacterTextSplitter(
        separators=[
            "\n\n",  # 空行
            "\n",  # 换行符
            " ",  # 空格
            ".",  # 英文句号
            ",",  # 英文逗号
            "\u3002",  # 中文句号
            "\uff0c",  # 中文逗号
            "\u3001",  # 中文顿号
            "",  # 空字符
        ],
        chunk_size=500,
        chunk_overlap=100,
        add_start_index=True
    )

    # 5. 对文档进行分块
    splits = text_splitter.split_documents(docs)

    # 6. 为每个分块添加source元数据
    for split in splits:
        split.metadata["source"] = file_path

    # 7. 打印前3个分块用于测试
    print(f"DOCX文件 {file_path} 被切分为 {len(splits)} 个分块。")
    for i, split in enumerate(splits[:3]):
        print(f"\n分块 {i + 1}:")
        print(split.page_content)
        print("元数据:", split.metadata)

    return splits

def load_pptx_simply(file_path:str)->List[Document]:
    # 负责人： 高哲文
    # 该方法基于python-pptx库加载pptx文件，只提取文本框和表格中的文本，并返回一个List[Document]
    try:
        # 1. 使用 python-pptx 加载 pptx 文件
        prs = Presentation(file_path)

        # 2. 提取每张幻灯片中的文本内容
        all_texts = []
        for slide in prs.slides:
            slide_text = ""

            # 提取文本框中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"

            # 提取表格中的文本
            for shape in slide.shapes:
                if shape.shape_type == 19:  # 19 对应表格
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            slide_text += cell.text.strip() + "\n"

            if slide_text.strip():
                all_texts.append(slide_text.strip())

        # 3. 使用 RecursiveCharacterTextSplitter 对提取的文本进行分割
        text_splitter = RecursiveCharacterTextSplitter(
            separators=[
                "\n\n",  # 空行，常作为不同章节的分割标志
                "\n",  # 换行符，常作为段落之间或文本与公式之间的分割标志
                " ",  # 空格，常作为英文单词之间分割标志
                ".",  # 句号
                ",",  # 逗号
                "\u200b",  # Zero-width space (零宽空格)
                "\uff0c",  # 全角逗号
                "\u3001",  # 顿号
                "\uff0e",  # 全角句号
                "\u3002",  # 句号
            ],
            chunk_size=500,  # 每个文档的最大字数
            chunk_overlap=200,  # 每个文档的重叠部分
            add_start_index=True
        )

        # 将所有幻灯片文本切割成多个 Document 对象
        all_docs = text_splitter.create_documents(all_texts)
        # 添加元数据source
        for doc in all_docs:
            doc.metadata["source"] = file_path

        # 4. 输出前三个 Document 对象，检查是否加载正确
        print(f"加载的 PPTX 文件包含 {len(all_docs)} 个文档块。\n")
        for i in range(len(all_docs)):
            print(f"Document {i + 1} 元数据: {all_docs[i].metadata}")
            print(f"Document {i + 1} 内容: {all_docs[i].page_content}\n")

        # 5. 返回所有的 Document 对象
        return all_docs

    except Exception as e:
        print(f"加载 PPTX 文件时出错: {e}")
        return []


def load_pdf_with_Azure(file_path:str)->List[Document]:
    # 该方法基于Azure Form Recognizer提取pdf文件中文本，返回一个List[Document]
    # 优点：文本提取精度高，可以准确提取 文本（包括手写）、表格、文档结构（例如标题、章节标题等）和 键值对
    # 缺点：付费服务(至少使用S0级别的订阅)

    # 1、检查变量endpoint和key是否设置了

    if  AZURE_ENDPOINT is None or AZURE_API_KEY is None :
        print("请设置环境变量AZURE_ENDPOINT和AZURE_API_KEY")
        return []

    loader = AzureAIDocumentIntelligenceLoader(
        api_endpoint=AZURE_ENDPOINT,
        api_key=AZURE_API_KEY,
        file_path=file_path,
        api_model="prebuilt-layout",
    )
    docs = loader.load()

    # 输出发现docs长度为1，唯一Document的内容为md格式的Document，元数据为大量无效信息。
    print("测试用的文件被加载为：", len(docs), "个Document\n")
    for i in range(len(docs)):
        print(f"第{i}页的内容为: {docs[i].page_content}\n\n")
        print(f"第{i}页的元数据为: {docs[i].metadata}\n")

    docs[0].metadata={"source":file_path}

    # 初始化 RecursiveCharacterTextSplitter
    text_splitter = RecursiveCharacterTextSplitter(
        separators=[
            "<!-- PageBreak -->",# 对于md格式的内容，先分页
            "##",    # 之后按照节标题切分
            "\n\n",  # 空行，常作为不同章节的分割标志
            "\n",  # 换行符，常作为段落之间或文本与公式之间的分割标志
            " ",  # 空格，常作为英文单词之间分割标志
            ".",
            ",",
            "\u200b",  # Zero-width space (零宽空格): （不可见）
            "\uff0c",  # Fullwidth comma (全角逗号): ，
            "\u3001",  # Ideographic comma (顿号): 、
            "\uff0e",  # Fullwidth full stop (全角句号): ．
            "\u3002",  # Ideographic full stop (句号): 。
        ],
        chunk_size=500,
        chunk_overlap=200,
        add_start_index=True
    )
    # 分割文档
    all_splits = text_splitter.split_documents(docs)
    # 输出前all_splits中前3个分块用于测试
    for i, split in enumerate(all_splits[:3]):
        print(f"\n分块 {i + 1}:")
        print(split.page_content)
        print("元数据:", split.metadata)
    return all_splits

if __name__ == "__main__":
    # load_txt("test_files/核工业百科.txt")
    # load_pdf_simply("test_files/1.10MW 高温堆热启动时蒸汽发生器.pdf")
    # load_md("test_files/LangChainItroduction.md")
    # load_docx("test_files/大创开题报告.docx")
    # load_pptx_simply("test_files/核工业专业知识问答模型构建-开题答辩.pptx")
     load_pdf_with_Azure("test_files/1.10MW 高温堆热启动时蒸汽发生器.pdf")